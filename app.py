import os
import tempfile
import zipfile
import subprocess
import shutil
import logging
from flask import Flask, request, send_file, jsonify
from flask_cors import CORS
from werkzeug.utils import secure_filename
import uuid
import re
from typing import Optional, Tuple, List, Dict, Any

# Try to import optional dependencies
try:
    import fitz  # PyMuPDF for PDF processing
    PYMUPDF_AVAILABLE = True
except ImportError:
    PYMUPDF_AVAILABLE = False
    fitz = None

try:
    from pptx import Presentation  # python-pptx for PPTX processing
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    Presentation = None

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

app = Flask(__name__)
CORS(app)
app.config['MAX_CONTENT_LENGTH'] = 50 * 1024 * 1024  # 50 MB limit

# Configuration
UPLOAD_FOLDER = 'uploads'
OUTPUT_FOLDER = 'output'
ALLOWED_EXTENSIONS = {
    'docx', 'doc', 'odt', 'rtf', 'html', 'htm', 'txt', 'md', 'markdown', 
    'tex', 'latex', 'epub', 'mobi', 'fb2', 'opml', 'org', 'mediawiki', 
    'dokuwiki', 'textile', 'rst', 'asciidoc', 'man', 'ms', 'docbook', 'xml',
    'jats', 'tei', 'ris', 'csljson', 'endnotexml', 'ipynb', 'csv', 'tsv',
    'json', 'native', 'typst', 'djot', 'creole', 'tikiwiki', 'twiki', 'vimwiki',
    'muse', 'pod', 't2t', 'haddock', 'mdoc', 'biblatex', 'bibtex', 'bits',
    'pdf', 'pptx'  # Added PDF and PPTX support
}

# Create directories if they don't exist
try:
    logger.info("Creating uploads directory...")
    os.makedirs(UPLOAD_FOLDER, exist_ok=True)
    logger.info("Creating output directory...")
    os.makedirs(OUTPUT_FOLDER, exist_ok=True)
    logger.info("Directories created successfully")
except Exception as e:
    logger.error(f"Error creating directories: {e}")
    raise

# Check Pandoc availability at startup
# try:
#     logger.info("Checking Pandoc availability...")
#     result = subprocess.run(['pandoc', '--version'], capture_output=True, text=True, timeout=10)
#     if result.returncode == 0:
#         logger.info(f"Pandoc is available: {result.stdout.split('\n')[0]}")
#     else:
#         logger.warning("Pandoc check failed")
# except Exception as e:
#     logger.warning(f"Could not check Pandoc: {e}")

logger.info("Application startup completed")

print("ALL ENV VARS:", dict(os.environ))

print("PORT ENV:", os.environ.get("PORT"))

def allowed_file(filename):
    return '.' in filename and \
           filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

def get_input_format(filename):
    """Determine input format based on file extension"""
    ext = filename.rsplit('.', 1)[1].lower()
    
    format_mapping = {
        'docx': 'docx',
        'doc': 'doc',
        'odt': 'odt',
        'rtf': 'rtf',
        'html': 'html',
        'htm': 'html',
        'txt': 'markdown',  # Use markdown for txt input
        'md': 'markdown',
        'markdown': 'markdown',
        'tex': 'latex',
        'latex': 'latex',
        'epub': 'epub',
        'mobi': 'mobi',
        'fb2': 'fb2',
        'opml': 'opml',
        'org': 'org',
        'mediawiki': 'mediawiki',
        'dokuwiki': 'dokuwiki',
        'textile': 'textile',
        'rst': 'rst',
        'asciidoc': 'asciidoc',
        'man': 'man',
        'ms': 'ms',
        'docbook': 'docbook',
        'xml': 'docbook',  # Default XML format
        'jats': 'jats',
        'tei': 'tei',
        'ris': 'ris',
        'csljson': 'csljson',
        'endnotexml': 'endnotexml',
        'ipynb': 'ipynb',
        'csv': 'csv',
        'tsv': 'tsv',
        'json': 'json',
        'native': 'native',
        'typst': 'typst',
        'djot': 'djot',
        'creole': 'creole',
        'tikiwiki': 'tikiwiki',
        'twiki': 'twiki',
        'vimwiki': 'vimwiki',
        'muse': 'muse',
        'pod': 'pod',
        't2t': 't2t',
        'haddock': 'haddock',
        'mdoc': 'mdoc',
        'biblatex': 'biblatex',
        'bibtex': 'bibtex',
        'bits': 'bits',
        'pdf': 'pdf',  # Added PDF support
        'pptx': 'pptx'  # Added PPTX support
    }
    
    return format_mapping.get(ext, 'markdown')

def get_format_suggestions(invalid_format):
    """Get suggestions for similar or common formats when an invalid format is entered"""
    common_formats = {
        'pdf': ['pdf'],
        'word': ['docx'],
        'powerpoint': ['pptx'],
        'html': ['html', 'html5', 'xhtml'],
        'markdown': ['markdown', 'gfm', 'commonmark'],
        'text': ['txt', 'plain'],
        'xml': ['xml', 'docbook', 'jats', 'tei'],
        'latex': ['latex', 'tex'],
        'epub': ['epub', 'epub2', 'epub3'],
        'presentation': ['revealjs', 'beamer', 's5', 'slidy'],
        'documentation': ['docbook', 'jats', 'asciidoc', 'rst']
    }
    
    suggestions = []
    invalid_lower = invalid_format.lower()
    
    for category, formats in common_formats.items():
        if invalid_lower in category or any(invalid_lower in fmt for fmt in formats):
            suggestions.extend(formats)
    
    # Add some common formats if no specific matches
    if not suggestions:
        suggestions = ['pdf', 'docx', 'html', 'markdown', 'txt', 'xml']
    
    return list(set(suggestions))[:5]  # Return up to 5 unique suggestions

def is_format_supported(output_format):
    """Check if the output format is supported by Pandoc - now more permissive"""
    # Allow any format input and let Pandoc handle validation
    # This makes the system more flexible for custom formats
    return True

def validate_output_file(output_path, output_format):
    """Validate that the output file was created correctly for the given format"""
    try:
        # Check if file exists and has content
        if not os.path.exists(output_path):
            return False, "Output file was not created"
        
        file_size = os.path.getsize(output_path)
        if file_size == 0:
            return False, "Output file is empty"
        
        # Basic validation for any format - just check if file has content
        # This allows maximum flexibility for any output format
        
        # For binary formats, just check file size
        if output_format in ['pdf', 'docx', 'pptx', 'odt', 'epub', 'mobi', 'fb2']:
            if file_size < 50:  # Even small binary files should have some content
                return False, f"Output {output_format} file appears to be corrupted (too small)"
        
        # For text-based formats, check for content
        else:
            try:
                with open(output_path, 'r', encoding='utf-8') as f:
                    content = f.read(100)  # Read first 100 chars
                    if not content.strip():
                        return False, f"Output {output_format} file is empty"
            except UnicodeDecodeError:
                # If it's not UTF-8, it might be a binary format we didn't expect
                # Just check if it has reasonable size
                if file_size < 50:
                    return False, f"Output {output_format} file appears to be corrupted (too small)"
        
        # If we get here, the file appears valid
        return True, f"Output {output_format} file validated successfully"
        
    except Exception as e:
        return False, f"Error validating output file: {str(e)}"

def convert_file_with_pandoc(input_path, output_path, input_format, output_format, extract_media_dir):
    """Convert file using Pandoc with precise format-specific options and enhanced media handling"""
    try:
        # Build base pandoc command
        cmd = [
            'pandoc',
            input_path,
            '-f', input_format,
            '-t', output_format,
            '-o', output_path
        ]
        
        # Add media extraction for ALL formats that might contain images
        # This ensures we capture images from any source format
        media_supporting_formats = [
            'html', 'html5', 'xhtml', 'epub', 'epub2', 'epub3', 'docx', 'pptx', 'odt',
            'markdown', 'gfm', 'commonmark', 'commonmark_x', 'rst', 'asciidoc',
            'textile', 'mediawiki', 'dokuwiki', 'org', 'opml', 'fb2', 'mobi',
            'docbook', 'docbook4', 'docbook5', 'jats', 'tei', 'icml'
        ]
        
        # Always extract media for formats that support it, or if input might contain images
        if output_format in media_supporting_formats or input_format in ['docx', 'pptx', 'odt', 'epub', 'html']:
            cmd.extend(['--extract-media', extract_media_dir])
        
        # Add format-specific options for precise conversion
        if output_format == 'gfm':
            cmd.extend(['--wrap=none', '--markdown-headings=atx'])
        elif output_format == 'markdown':
            cmd.extend(['--wrap=none'])
        elif output_format == 'html':
            cmd.extend(['--standalone', '--self-contained'])
        elif output_format == 'html5':
            cmd.extend(['--standalone', '--self-contained', '--to=html5'])
        elif output_format == 'xhtml':
            cmd.extend(['--standalone', '--self-contained', '--to=xhtml'])
        elif output_format == 'pdf':
            cmd.extend(['--pdf-engine=xelatex'])
        elif output_format == 'latex':
            cmd.extend(['--standalone'])
        elif output_format == 'docx':
            cmd.extend(['--reference-doc='])  # Use default template
        elif output_format == 'pptx':
            cmd.extend(['--reference-doc='])  # Use default template
        elif output_format == 'odt':
            cmd.extend(['--reference-doc='])  # Use default template
        elif output_format == 'rtf':
            cmd.extend([])  # No special options needed
        elif output_format == 'epub':
            cmd.extend(['--epub-cover-image='])  # No cover image
        elif output_format == 'epub2':
            cmd.extend(['--to=epub2'])
        elif output_format == 'epub3':
            cmd.extend(['--to=epub3'])
        elif output_format == 'txt':
            cmd.extend(['--wrap=none'])
        elif output_format == 'xml':
            cmd.extend(['--standalone'])
        elif output_format == 'docbook':
            cmd.extend(['--standalone', '--to=docbook5'])
        elif output_format == 'docbook5':
            cmd.extend(['--standalone', '--to=docbook5'])
        elif output_format == 'docbook4':
            cmd.extend(['--standalone', '--to=docbook4'])
        elif output_format == 'jats':
            cmd.extend(['--standalone', '--to=jats'])
        elif output_format == 'jats_archiving':
            cmd.extend(['--standalone', '--to=jats_archiving'])
        elif output_format == 'jats_publishing':
            cmd.extend(['--standalone', '--to=jats_publishing'])
        elif output_format == 'jats_articleauthoring':
            cmd.extend(['--standalone', '--to=jats_articleauthoring'])
        elif output_format == 'revealjs':
            cmd.extend(['--standalone', '--to=revealjs'])
        elif output_format == 'beamer':
            cmd.extend(['--pdf-engine=xelatex', '--to=beamer'])
        elif output_format == 's5':
            cmd.extend(['--standalone', '--to=s5'])
        elif output_format == 'slideous':
            cmd.extend(['--standalone', '--to=slideous'])
        elif output_format == 'dzslides':
            cmd.extend(['--standalone', '--to=dzslides'])
        elif output_format == 'slidy':
            cmd.extend(['--standalone', '--to=slidy'])
        elif output_format == 'asciidoc':
            cmd.extend(['--wrap=none'])
        elif output_format == 'rst':
            cmd.extend(['--wrap=none'])
        elif output_format == 'org':
            cmd.extend(['--wrap=none'])
        elif output_format == 'textile':
            cmd.extend(['--wrap=none'])
        elif output_format == 'mediawiki':
            cmd.extend(['--wrap=none'])
        elif output_format == 'dokuwiki':
            cmd.extend(['--wrap=none'])
        elif output_format == 'haddock':
            cmd.extend(['--wrap=none'])
        elif output_format == 'man':
            cmd.extend([])
        elif output_format == 'ms':
            cmd.extend([])
        elif output_format == 'opml':
            cmd.extend(['--standalone'])
        elif output_format == 'fb2':
            cmd.extend(['--standalone'])
        elif output_format == 'mobi':
            cmd.extend(['--standalone'])
        elif output_format == 'icml':
            cmd.extend(['--standalone'])
        elif output_format == 'tei':
            cmd.extend(['--standalone'])
        elif output_format == 'native':
            cmd.extend([])
        elif output_format == 'json':
            cmd.extend(['--to=json'])
        elif output_format == 'commonmark':
            cmd.extend(['--wrap=none', '--to=commonmark'])
        elif output_format == 'commonmark_x':
            cmd.extend(['--wrap=none', '--to=commonmark_x'])
        elif output_format == 'markua':
            cmd.extend(['--wrap=none', '--to=markua'])
        elif output_format == 'spip':
            cmd.extend(['--wrap=none'])
        elif output_format == 'texinfo':
            cmd.extend(['--standalone'])
        elif output_format == 'opendocument':
            cmd.extend(['--to=opendocument'])
        else:
            # For any other format, try direct conversion without special options
            # This allows maximum flexibility for custom formats
            cmd.extend([])
        
        # Log the command being executed (without sensitive info)
        logger.info(f"Executing pandoc command: {' '.join(cmd[:4])} ... [output format: {output_format}]")
        
        # Execute pandoc command
        result = subprocess.run(cmd, capture_output=True, text=True, check=True)
        
        # Check if media was extracted
        media_files = []
        if os.path.exists(extract_media_dir):
            for root, dirs, files in os.walk(extract_media_dir):
                for file in files:
                    media_files.append(os.path.join(root, file))
        
        if media_files:
            logger.info(f"Extracted {len(media_files)} media files to {extract_media_dir}")
        else:
            logger.info("No media files were extracted")
        
        # Validate the output file
        validation_success, validation_message = validate_output_file(output_path, output_format)
        if not validation_success:
            return False, f"Conversion completed but validation failed: {validation_message}"
        
        return True, None
        
    except subprocess.CalledProcessError as e:
        error_msg = f"Pandoc error: {e.stderr}"
        logger.error(error_msg)
        return False, error_msg
    except FileNotFoundError:
        error_msg = "Pandoc is not installed or not found in PATH"
        logger.error(error_msg)
        return False, error_msg
    except Exception as e:
        error_msg = f"Conversion error: {str(e)}"
        logger.error(error_msg)
        return False, error_msg

def convert_pdf_to_markdown(pdf_path: str, output_path: str) -> Tuple[bool, Optional[str]]:
    """Convert PDF to Markdown using PyMuPDF"""
    if not PYMUPDF_AVAILABLE:
        return False, "PyMuPDF is not available for PDF processing"
    
    try:
        doc = fitz.open(pdf_path)
        markdown_content = []
        
        for page_num in range(len(doc)):
            page = doc.load_page(page_num)
            text = page.get_text()
            
            # Clean up the text
            lines = text.split('\n')
            cleaned_lines = []
            
            for line in lines:
                line = line.strip()
                if line:
                    # Try to detect headers (simple heuristic)
                    if len(line) < 100 and line.isupper():
                        cleaned_lines.append(f"# {line.title()}")
                    else:
                        cleaned_lines.append(line)
            
            if cleaned_lines:
                markdown_content.extend(cleaned_lines)
                markdown_content.append('')  # Add blank line between pages
        
        doc.close()
        
        # Write to markdown file
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write('\n'.join(markdown_content))
        
        return True, None
        
    except Exception as e:
        return False, f"PDF conversion error: {str(e)}"

def convert_pptx_to_markdown(pptx_path: str, output_path: str) -> Tuple[bool, Optional[str]]:
    """Convert PPTX to Markdown using python-pptx"""
    if not PPTX_AVAILABLE:
        return False, "python-pptx is not available for PPTX processing"
    
    try:
        prs = Presentation(pptx_path)
        markdown_content = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            # Add slide header
            markdown_content.append(f"# Slide {slide_num}")
            markdown_content.append('')
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text = shape.text.strip()
                    
                    # Try to detect if it's a title or content
                    if len(text) < 100 and text.isupper():
                        markdown_content.append(f"## {text.title()}")
                    else:
                        # Split into paragraphs
                        paragraphs = text.split('\n')
                        for para in paragraphs:
                            if para.strip():
                                markdown_content.append(para.strip())
                                markdown_content.append('')
            
            markdown_content.append('---')  # Slide separator
            markdown_content.append('')
        
        # Write to markdown file
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write('\n'.join(markdown_content))
        
        # Ensure the markdown file was created
        if os.path.exists(output_path):
            return True, None
        else:
            return False, "Markdown file was not created from PPTX"
        
    except Exception as e:
        return False, f"PPTX conversion error: {str(e)}"

def preprocess_special_formats(input_path, input_format, temp_dir):
    """Preprocess special formats (PDF, PPTX) to convert them to Pandoc-supported formats. PPTX and others are handled directly by Pandoc."""
    try:
        if input_format == 'pdf':
            # Convert PDF to Markdown first
            temp_md_path = os.path.join(temp_dir, 'temp_converted.md')
            success, error = convert_pdf_to_markdown(input_path, temp_md_path)
            if success and os.path.exists(temp_md_path):
                return temp_md_path, 'markdown'
            else:
                return None, error or "PDF to Markdown conversion failed"
        elif input_format == 'pptx':
            # Convert PPTX to Markdown first
            temp_md_path = os.path.join(temp_dir, 'temp_converted.md')
            success, error = convert_pptx_to_markdown(input_path, temp_md_path)
            if success and os.path.exists(temp_md_path):
                return temp_md_path, 'markdown'
            else:
                return None, error or "PPTX to Markdown conversion failed"
        else:
            # No preprocessing needed for other formats (including PPTX)
            return input_path, input_format
    except Exception as e:
        return None, f"Preprocessing error: {str(e)}"

def map_output_format(user_format):
    """Map user-friendly format names to actual Pandoc format names"""
    format_mapping = {
        'txt': 'plain',  # Pandoc uses 'plain' for plain text, not 'txt'
        'text': 'plain',
        'plaintext': 'plain',
        'word': 'docx',
        'powerpoint': 'pptx',
        'presentation': 'pptx',
        'document': 'docx',
        'webpage': 'html',
        'web': 'html',
        'page': 'html',
        'notebook': 'ipynb',
        'jupyter': 'ipynb',
        'ebook': 'epub',
        'book': 'epub',
        'slide': 'revealjs',
        'slides': 'revealjs',
        'deck': 'revealjs',
        'beamer_slide': 'beamer',
        'latex_slide': 'beamer',
        'pdf_slide': 'beamer',
        'markdown_github': 'gfm',
        'github_markdown': 'gfm',
        'github': 'gfm',
        'commonmark_x': 'commonmark_x',
        'commonmark_x_extended': 'commonmark_x',
        'extended_commonmark': 'commonmark_x',
        'markua': 'markua',
        'markua_document': 'markua',
        'spip': 'spip',
        'spip_wiki': 'spip',
        'epub2': 'epub2',
        'epub_version2': 'epub2',
        'epub3': 'epub3',
        'epub_version3': 'epub3',
        'docbook4': 'docbook4',
        'docbook_version4': 'docbook4',
        'docbook5': 'docbook5',
        'docbook_version5': 'docbook5',
        'jats_archiving': 'jats_archiving',
        'jats_archiving_version': 'jats_archiving',
        'jats_publishing': 'jats_publishing',
        'jats_publishing_version': 'jats_publishing',
        'jats_articleauthoring': 'jats_articleauthoring',
        'jats_article_authoring': 'jats_articleauthoring',
        'html5': 'html5',
        'html_version5': 'html5',
        'html4': 'html4',
        'html_version4': 'html4',
        'xhtml': 'xhtml',
        'xhtml5': 'xhtml5',
        'xhtml_version5': 'xhtml5',
        'xhtml4': 'xhtml4',
        'xhtml_version4': 'xhtml4',
        'markdown_github': 'markdown_github',
        'markdown_mmd': 'markdown_mmd',
        'markdown_phpextra': 'markdown_phpextra',
        'markdown_strict': 'markdown_strict',
        'markdown_texinfo': 'markdown_texinfo',
        'commonmark': 'commonmark',
        'commonmark_strict': 'commonmark',
        'commonmark_x': 'commonmark_x',
        'commonmark_extended': 'commonmark_x',
        'gfm': 'gfm',
        'github_flavored': 'gfm',
        'markua': 'markua',
        'markua_document': 'markua',
        'spip': 'spip',
        'spip_wiki': 'spip',
        'epub2': 'epub2',
        'epub_version2': 'epub2',
        'epub3': 'epub3',
        'epub_version3': 'epub3',
        'texinfo': 'texinfo',
        'tex_info': 'texinfo',
        'textile': 'textile',
        'textile_wiki': 'textile',
        'org': 'org',
        'org_mode': 'org',
        'emacs_org': 'org',
        'asciidoc': 'asciidoc',
        'ascii_doc': 'asciidoc',
        'rst': 'rst',
        'restructuredtext': 'rst',
        'rest': 'rst',
        'mediawiki': 'mediawiki',
        'wiki': 'mediawiki',
        'wikipedia': 'mediawiki',
        'dokuwiki': 'dokuwiki',
        'doku_wiki': 'dokuwiki',
        'haddock': 'haddock',
        'haskell_doc': 'haddock',
        'opml': 'opml',
        'outline': 'opml',
        'fb2': 'fb2',
        'fictionbook': 'fb2',
        'mobi': 'mobi',
        'kindle': 'mobi',
        'icml': 'icml',
        'indesign': 'icml',
        'tei': 'tei',
        'text_encoding_initiative': 'tei',
        'native': 'native',
        'pandoc_native': 'native',
        'json': 'json',
        'javascript_object_notation': 'json',
        'jats_archiving': 'jats_archiving',
        'jats_archiving_version': 'jats_archiving',
        'jats_publishing': 'jats_publishing',
        'jats_publishing_version': 'jats_publishing',
        'jats_articleauthoring': 'jats_articleauthoring',
        'jats_article_authoring': 'jats_articleauthoring',
        'html5': 'html5',
        'html_version5': 'html5',
        'html4': 'html4',
        'html_version4': 'html4',
        'xhtml': 'xhtml',
        'xhtml5': 'xhtml5',
        'xhtml_version5': 'xhtml5',
        'xhtml4': 'xhtml4',
        'xhtml_version4': 'xhtml4',
        'markdown_github': 'markdown_github',
        'markdown_mmd': 'markdown_mmd',
        'markdown_phpextra': 'markdown_phpextra',
        'markdown_strict': 'markdown_strict',
        'markdown_texinfo': 'markdown_texinfo',
        'commonmark': 'commonmark',
        'commonmark_x': 'commonmark_x',
        'gfm': 'gfm',
        'markua': 'markua',
        'spip': 'txt',
        'epub2': 'epub',
        'epub3': 'epub',
        'docbook4': 'xml',
        'docbook5': 'xml',
        'man': 'man',
        'ms': 'ms',
        'texinfo': 'texi',
        'textile': 'textile',
        'org': 'org',
        'asciidoc': 'adoc',
        'rst': 'rst',
        'mediawiki': 'wiki',
        'dokuwiki': 'txt',
        'haddock': 'hs',
        'opml': 'opml',
        'fb2': 'fb2',
        'mobi': 'mobi',
        'icml': 'icml',
        'tei': 'xml',
        'native': 'native',
        'json': 'json',
        'jats_archiving': 'xml',
        'jats_publishing': 'xml',
        'jats_articleauthoring': 'xml',
        'html5': 'html',
        'html4': 'html',
        'xhtml': 'xhtml',
        'xhtml5': 'xhtml',
        'xhtml4': 'xhtml',
        'markdown_github': 'md',
        'markdown_mmd': 'md',
        'markdown_phpextra': 'md',
        'markdown_strict': 'md',
        'markdown_texinfo': 'texi',
        'commonmark': 'md',
        'commonmark_x': 'md',
        'gfm': 'md',
        'markua': 'md',
        'spip': 'txt',
        'epub2': 'epub',
        'epub3': 'epub',
        'texinfo': 'texi'
    }
    
    # Return the mapped format or the original if no mapping exists
    return format_mapping.get(user_format.lower(), user_format.lower())

def fix_image_paths_in_file(file_path: str, media_dir: str, output_format: str) -> bool:
    """Fix image paths in converted files to use relative paths to img/ folder"""
    try:
        # Read the file content
        with open(file_path, 'r', encoding='utf-8') as f:
            content = f.read()
        
        # Define patterns for different output formats
        if output_format in ['html', 'html5', 'xhtml']:
            # Fix HTML img src attributes
            content = re.sub(
                r'src="([^"]*)"',
                lambda m: f'src="img/{os.path.basename(m.group(1))}"' if m.group(1) else m.group(0),
                content
            )
        elif output_format in ['markdown', 'gfm', 'commonmark', 'commonmark_x']:
            # Fix Markdown image syntax
            content = re.sub(
                r'!\[([^\]]*)\]\(([^)]+)\)',
                lambda m: f'![{m.group(1)}](img/{os.path.basename(m.group(2))})',
                content
            )
        elif output_format in ['rst']:
            # Fix reStructuredText image syntax
            content = re.sub(
                r'\.\. image:: ([^\n]+)',
                lambda m: f'.. image:: img/{os.path.basename(m.group(1))}',
                content
            )
        elif output_format in ['asciidoc']:
            # Fix AsciiDoc image syntax
            content = re.sub(
                r'image::([^[]+)\[([^\]]*)\]',
                lambda m: f'image::img/{os.path.basename(m.group(1))}[{m.group(2)}]',
                content
            )
        
        # Write the fixed content back
        with open(file_path, 'w', encoding='utf-8') as f:
            f.write(content)
        
        return True
        
    except Exception as e:
        logger.error(f"Error fixing image paths in {file_path}: {str(e)}")
        return False

def organize_media_files(media_dir: str, img_dir: str) -> List[str]:
    """Organize extracted media files into img/ folder and return list of moved files"""
    moved_files = []
    
    try:
        # Create img directory if it doesn't exist
        os.makedirs(img_dir, exist_ok=True)
        
        # Move all media files to img/ folder
        if os.path.exists(media_dir):
            for root, dirs, files in os.walk(media_dir):
                for file in files:
                    src_path = os.path.join(root, file)
                    dst_path = os.path.join(img_dir, file)
                    
                    # Handle duplicate filenames
                    counter = 1
                    base_name, ext = os.path.splitext(file)
                    while os.path.exists(dst_path):
                        new_name = f"{base_name}_{counter}{ext}"
                        dst_path = os.path.join(img_dir, new_name)
                        counter += 1
                    
                    # Move the file
                    shutil.move(src_path, dst_path)
                    moved_files.append(dst_path)
                    logger.info(f"Moved media file: {file} -> img/{os.path.basename(dst_path)}")
        
        return moved_files
        
    except Exception as e:
        logger.error(f"Error organizing media files: {str(e)}")
        return moved_files

@app.route('/convert', methods=['POST'])
def convert_files():
    try:
        if 'files' not in request.files:
            return jsonify({'error': 'No files provided'}), 400
        
        files = request.files.getlist('files')
        original_output_format = request.form.get('output_format', 'pdf')
        output_format = original_output_format.strip().lower()
        
        if not files or all(file.filename == '' for file in files):
            return jsonify({'error': 'No files selected'}), 400
        
        # Validate output format
        if not output_format or output_format.strip() == '':
            return jsonify({'error': 'Output format is required'}), 400
        
        # Map user-friendly format names to actual Pandoc format names
        output_format = map_output_format(output_format)
        
        logger.info(f"User requested format: {original_output_format}, mapped to: {output_format}")
        
        # Create unique session directory
        session_id = str(uuid.uuid4())
        session_dir = os.path.join(OUTPUT_FOLDER, session_id)
        os.makedirs(session_dir, exist_ok=True)
        
        # Create subdirectories
        uploads_dir = os.path.join(session_dir, 'uploads')
        converted_dir = os.path.join(session_dir, 'converted')
        media_dir = os.path.join(session_dir, 'media')
        img_dir = os.path.join(session_dir, 'img')  # New img directory for organized images
        
        os.makedirs(uploads_dir, exist_ok=True)
        os.makedirs(converted_dir, exist_ok=True)
        os.makedirs(media_dir, exist_ok=True)
        os.makedirs(img_dir, exist_ok=True)
        
        conversion_errors = []
        converted_files = []
        
        for file in files:
            if file and file.filename and allowed_file(file.filename):
                # Save uploaded file
                filename = secure_filename(file.filename)
                input_path = os.path.join(uploads_dir, filename)
                file.save(input_path)
                logger.info(f"Saved uploaded file to: {input_path} (exists: {os.path.exists(input_path)})")
                
                # Determine input format
                input_format = get_input_format(filename)
                logger.info(f"Detected input format for {filename}: {input_format}")
                
                # Preprocess special formats (PDF, PPTX) if needed
                processed_input_path, processed_input_format = preprocess_special_formats(
                    input_path, input_format, uploads_dir
                )
                logger.info(f"Processed input path: {processed_input_path} (exists: {os.path.exists(processed_input_path) if processed_input_path else 'N/A'})")
                
                if processed_input_path is None:
                    logger.error(f"Failed to preprocess {filename}: {processed_input_format}")
                    conversion_errors.append(f"{filename}: {processed_input_format}")
                    continue
                
                # Check if processed input file exists before Pandoc call
                if not os.path.exists(processed_input_path):
                    logger.error(f"File does not exist before Pandoc call: {processed_input_path}")
                    conversion_errors.append(f"{filename}: Input file missing before Pandoc call")
                    continue
                else:
                    logger.info(f"File confirmed to exist before Pandoc call: {processed_input_path}")
                    logger.info(f"Current working directory: {os.getcwd()}")
                
                # Generate output filename
                base_name = os.path.splitext(filename)[0]
                
                # Define common format extensions
                format_extensions = {
                    'gfm': 'md',
                    'markdown': 'md',
                    'html': 'html',
                    'latex': 'tex',
                    'pdf': 'pdf',
                    'docx': 'docx',
                    'odt': 'odt',
                    'rtf': 'rtf',
                    'epub': 'epub',
                    'pptx': 'pptx',
                    'xml': 'xml',
                    'txt': 'txt',
                    'plain': 'txt',
                    'docbook': 'xml',
                    'jats': 'xml',
                    'opendocument': 'odt',
                    'revealjs': 'html',
                    'beamer': 'pdf',
                    's5': 'html',
                    'slideous': 'html',
                    'dzslides': 'html',
                    'slidy': 'html',
                    'asciidoc': 'adoc',
                    'rst': 'rst',
                    'org': 'org',
                    'textile': 'textile',
                    'mediawiki': 'wiki',
                    'dokuwiki': 'txt',
                    'haddock': 'hs',
                    'man': 'man',
                    'ms': 'ms',
                    'opml': 'opml',
                    'fb2': 'fb2',
                    'mobi': 'mobi',
                    'icml': 'icml',
                    'tei': 'xml',
                    'native': 'native',
                    'json': 'json',
                    'docbook5': 'xml',
                    'docbook4': 'xml',
                    'jats_archiving': 'xml',
                    'jats_publishing': 'xml',
                    'jats_articleauthoring': 'xml',
                    'html5': 'html',
                    'html4': 'html',
                    'xhtml': 'xhtml',
                    'xhtml5': 'xhtml',
                    'xhtml4': 'xhtml',
                    'markdown_github': 'md',
                    'markdown_mmd': 'md',
                    'markdown_phpextra': 'md',
                    'markdown_strict': 'md',
                    'markdown_texinfo': 'texi',
                    'commonmark': 'md',
                    'commonmark_x': 'md',
                    'gfm': 'md',
                    'markua': 'md',
                    'spip': 'txt',
                    'epub2': 'epub',
                    'epub3': 'epub',
                    'docbook4': 'xml',
                    'docbook5': 'xml',
                    'man': 'man',
                    'ms': 'ms',
                    'texinfo': 'texi',
                    'textile': 'textile',
                    'org': 'org',
                    'asciidoc': 'adoc',
                    'rst': 'rst',
                    'mediawiki': 'wiki',
                    'dokuwiki': 'txt',
                    'haddock': 'hs',
                    'opml': 'opml',
                    'fb2': 'fb2',
                    'mobi': 'mobi',
                    'icml': 'icml',
                    'tei': 'xml',
                    'native': 'native',
                    'json': 'json',
                    'jats_archiving': 'xml',
                    'jats_publishing': 'xml',
                    'jats_articleauthoring': 'xml',
                    'html5': 'html',
                    'html4': 'html',
                    'xhtml': 'xhtml',
                    'xhtml5': 'xhtml',
                    'xhtml4': 'xhtml',
                    'markdown_github': 'md',
                    'markdown_mmd': 'md',
                    'markdown_phpextra': 'md',
                    'markdown_strict': 'md',
                    'markdown_texinfo': 'texi',
                    'commonmark': 'md',
                    'commonmark_x': 'md',
                    'gfm': 'md',
                    'markua': 'md',
                    'spip': 'txt',
                    'epub2': 'epub',
                    'epub3': 'epub',
                    'texinfo': 'texi'
                }
                
                # Get extension for the output format
                if output_format is None:
                    extension = 'txt'  # Default fallback
                else:
                    extension = format_extensions.get(output_format, output_format)
                output_filename = f"{base_name}.{extension}"
                
                output_path = os.path.join(converted_dir, output_filename)
                
                # Convert file
                logger.info(f"Converting {filename} from {input_format} to {output_format}")
                success, error = convert_file_with_pandoc(
                    processed_input_path, output_path, processed_input_format, output_format, media_dir
                )
                
                if success:
                    # Organize media files into img/ folder
                    moved_media_files = organize_media_files(media_dir, img_dir)
                    
                    # Fix image paths in the converted file if it's a text-based format
                    text_based_formats = ['html', 'html5', 'xhtml', 'markdown', 'gfm', 'commonmark', 
                                         'commonmark_x', 'rst', 'asciidoc', 'textile', 'mediawiki', 
                                         'dokuwiki', 'org', 'opml', 'fb2', 'mobi', 'docbook', 
                                         'docbook4', 'docbook5', 'jats', 'tei', 'icml']
                    
                    if output_format in text_based_formats and moved_media_files:
                        fix_image_paths_in_file(output_path, img_dir, output_format)
                        logger.info(f"Fixed image paths in {output_filename}")
                    
                    logger.info(f"Successfully converted {filename} to {output_filename} and validated output")
                    converted_files.append(output_filename)
                else:
                    logger.error(f"Failed to convert {filename}: {error}")
                    conversion_errors.append(f"{filename}: {error}")
            else:
                conversion_errors.append(f"Invalid file type: {file.filename}")
        
        if not converted_files:
            error_msg = "No files were successfully converted."
            if conversion_errors:
                error_msg += " Errors: " + "; ".join(conversion_errors)
            return jsonify({'error': error_msg}), 400
        
        # Create ZIP file
        zip_filename = f"converted_files_{session_id}.zip"
        zip_path = os.path.join(session_dir, zip_filename)
        
        with zipfile.ZipFile(zip_path, 'w', zipfile.ZIP_DEFLATED) as zipf:
            # Add converted files
            converted_count = 0
            for filename in os.listdir(converted_dir):
                file_path = os.path.join(converted_dir, filename)
                zipf.write(file_path, filename)
                converted_count += 1
                logger.info(f"Added converted file to ZIP: {filename}")
            
            # Add img folder with organized media files
            img_count = 0
            if os.path.exists(img_dir) and os.listdir(img_dir):
                for root, dirs, files in os.walk(img_dir):
                    for file in files:
                        file_path = os.path.join(root, file)
                        # Preserve img/ directory structure in ZIP
                        arcname = os.path.relpath(file_path, session_dir)
                        zipf.write(file_path, arcname)
                        img_count += 1
                        logger.info(f"Added image file to ZIP: {arcname}")
            
            logger.info(f"ZIP created with {converted_count} converted files and {img_count} image files")
        
        # Clean up uploaded files (keep converted and zip for download)
        shutil.rmtree(uploads_dir, ignore_errors=True)
        
        # Return ZIP file
        return send_file(
            zip_path,
            as_attachment=True,
            download_name=zip_filename,
            mimetype='application/zip'
        )
        
    except Exception as e:
        return jsonify({'error': f'Server error: {str(e)}'}), 500

@app.route('/test', methods=['GET'])
def test_conversion():
    """Test endpoint to verify Pandoc is working and check supported formats"""
    try:
        # Test Pandoc availability
        result = subprocess.run(['pandoc', '--version'], capture_output=True, text=True, timeout=10)
        if result.returncode == 0:
            pandoc_version = result.stdout.split('\n')[0]
        else:
            return jsonify({'error': 'Pandoc not available'}), 500
        
        # Test supported formats
        result = subprocess.run(['pandoc', '--list-input-formats'], capture_output=True, text=True, timeout=10)
        input_formats = result.stdout.strip().split('\n') if result.returncode == 0 else []
        
        result = subprocess.run(['pandoc', '--list-output-formats'], capture_output=True, text=True, timeout=10)
        output_formats = result.stdout.strip().split('\n') if result.returncode == 0 else []
        
        return jsonify({
            'pandoc_version': pandoc_version,
            'input_formats': input_formats,
            'output_formats': output_formats,
            'status': 'Pandoc is working correctly'
        }), 200
        
    except Exception as e:
        return jsonify({'error': f'Test failed: {str(e)}'}), 500

@app.route('/health', methods=['GET'])
def health_check():
    return "OK", 200

@app.route('/retry', methods=['POST'])
def retry_conversion():
    """Retry a failed conversion with the same parameters"""
    try:
        if 'files' not in request.files:
            return jsonify({'error': 'No files provided for retry'}), 400
        
        files = request.files.getlist('files')
        original_output_format = request.form.get('output_format', 'pdf')
        output_format = original_output_format.strip().lower()
        
        if not files or all(file.filename == '' for file in files):
            return jsonify({'error': 'No files selected for retry'}), 400
        
        # Validate output format
        if not output_format or output_format.strip() == '':
            return jsonify({'error': 'Output format is required for retry'}), 400
        
        # Map user-friendly format names to actual Pandoc format names
        output_format = map_output_format(output_format)
        
        logger.info(f"User requested format: {original_output_format}, mapped to: {output_format}")
        
        # Log retry attempt
        logger.info(f"Retry conversion requested for {len(files)} files to {output_format}")
        
        # Use the same conversion logic as /convert
        return convert_files()
        
    except Exception as e:
        logger.error(f"Retry conversion error: {str(e)}")
        return jsonify({'error': f'Retry failed: {str(e)}'}), 500

@app.route('/')
def index():
    return "Backend server is running. Use /convert, /retry, or /health endpoints."

def get_supported_output_formats(input_formats):
    """Get list of supported output formats for given input formats"""
    
    # Define format compatibility matrix
    # This maps input formats to their supported output formats
    format_compatibility = {
        # Document formats
        'docx': [
            'html', 'html5', 'xhtml', 'markdown', 'gfm', 'commonmark', 'commonmark_x',
            'pdf', 'latex', 'docbook', 'docbook4', 'docbook5', 'jats', 'tei',
            'epub', 'epub2', 'epub3', 'mobi', 'fb2', 'rtf', 'txt', 'plain',
            'json', 'native', 'icml', 'opml', 'org', 'textile', 'mediawiki',
            'dokuwiki', 'haddock', 'man', 'ms', 'asciidoc', 'rst', 'opendocument'
        ],
        'doc': [
            'html', 'html5', 'xhtml', 'markdown', 'gfm', 'commonmark', 'commonmark_x',
            'pdf', 'latex', 'docbook', 'docbook4', 'docbook5', 'jats', 'tei',
            'epub', 'epub2', 'epub3', 'mobi', 'fb2', 'rtf', 'txt', 'plain',
            'json', 'native', 'icml', 'opml', 'org', 'textile', 'mediawiki',
            'dokuwiki', 'haddock', 'man', 'ms', 'asciidoc', 'rst', 'opendocument'
        ],
        'odt': [
            'html', 'html5', 'xhtml', 'markdown', 'gfm', 'commonmark', 'commonmark_x',
            'pdf', 'latex', 'docbook', 'docbook4', 'docbook5', 'jats', 'tei',
            'epub', 'epub2', 'epub3', 'mobi', 'fb2', 'rtf', 'txt', 'plain',
            'json', 'native', 'icml', 'opml', 'org', 'textile', 'mediawiki',
            'dokuwiki', 'haddock', 'man', 'ms', 'asciidoc', 'rst', 'docx'
        ],
        'rtf': [
            'html', 'html5', 'xhtml', 'markdown', 'gfm', 'commonmark', 'commonmark_x',
            'pdf', 'latex', 'docbook', 'docbook4', 'docbook5', 'jats', 'tei',
            'epub', 'epub2', 'epub3', 'mobi', 'fb2', 'txt', 'plain',
            'json', 'native', 'icml', 'opml', 'org', 'textile', 'mediawiki',
            'dokuwiki', 'haddock', 'man', 'ms', 'asciidoc', 'rst', 'docx', 'odt'
        ],
        
        # Presentation formats
        'pptx': [
            'html', 'html5', 'xhtml', 'markdown', 'gfm', 'commonmark', 'commonmark_x',
            'pdf', 'latex', 'docbook', 'docbook4', 'docbook5', 'jats', 'tei',
            'epub', 'epub2', 'epub3', 'mobi', 'fb2', 'rtf', 'txt', 'plain',
            'json', 'native', 'icml', 'opml', 'org', 'textile', 'mediawiki',
            'dokuwiki', 'haddock', 'man', 'ms', 'asciidoc', 'rst', 'revealjs',
            'beamer', 's5', 'slideous', 'dzslides', 'slidy'
        ],
        
        # Web formats
        'html': [
            'markdown', 'gfm', 'commonmark', 'commonmark_x', 'pdf', 'latex',
            'docbook', 'docbook4', 'docbook5', 'jats', 'tei', 'epub', 'epub2',
            'epub3', 'mobi', 'fb2', 'rtf', 'txt', 'plain', 'json', 'native',
            'icml', 'opml', 'org', 'textile', 'mediawiki', 'dokuwiki', 'haddock',
            'man', 'ms', 'asciidoc', 'rst', 'html5', 'xhtml', 'docx', 'odt'
        ],
        'htm': [
            'markdown', 'gfm', 'commonmark', 'commonmark_x', 'pdf', 'latex',
            'docbook', 'docbook4', 'docbook5', 'jats', 'tei', 'epub', 'epub2',
            'epub3', 'mobi', 'fb2', 'rtf', 'txt', 'plain', 'json', 'native',
            'icml', 'opml', 'org', 'textile', 'mediawiki', 'dokuwiki', 'haddock',
            'man', 'ms', 'asciidoc', 'rst', 'html5', 'xhtml', 'docx', 'odt'
        ],
        
        # Markup formats
        'markdown': [
            'html', 'html5', 'xhtml', 'gfm', 'commonmark', 'commonmark_x',
            'pdf', 'latex', 'docbook', 'docbook4', 'docbook5', 'jats', 'tei',
            'epub', 'epub2', 'epub3', 'mobi', 'fb2', 'rtf', 'txt', 'plain',
            'json', 'native', 'icml', 'opml', 'org', 'textile', 'mediawiki',
            'dokuwiki', 'haddock', 'man', 'ms', 'asciidoc', 'rst', 'docx', 'odt',
            'revealjs', 'beamer', 's5', 'slideous', 'dzslides', 'slidy'
        ],
        'md': [
            'html', 'html5', 'xhtml', 'gfm', 'commonmark', 'commonmark_x',
            'pdf', 'latex', 'docbook', 'docbook4', 'docbook5', 'jats', 'tei',
            'epub', 'epub2', 'epub3', 'mobi', 'fb2', 'rtf', 'txt', 'plain',
            'json', 'native', 'icml', 'opml', 'org', 'textile', 'mediawiki',
            'dokuwiki', 'haddock', 'man', 'ms', 'asciidoc', 'rst', 'docx', 'odt',
            'revealjs', 'beamer', 's5', 'slideous', 'dzslides', 'slidy'
        ],
        
        # PDF (limited support)
        'pdf': [
            'markdown', 'gfm', 'commonmark', 'commonmark_x', 'txt', 'plain',
            'html', 'html5', 'xhtml', 'json', 'native'
        ],
        
        # E-book formats
        'epub': [
            'html', 'html5', 'xhtml', 'markdown', 'gfm', 'commonmark', 'commonmark_x',
            'pdf', 'latex', 'docbook', 'docbook4', 'docbook5', 'jats', 'tei',
            'mobi', 'fb2', 'rtf', 'txt', 'plain', 'json', 'native', 'icml',
            'opml', 'org', 'textile', 'mediawiki', 'dokuwiki', 'haddock', 'man',
            'ms', 'asciidoc', 'rst', 'docx', 'odt'
        ],
        'mobi': [
            'html', 'html5', 'xhtml', 'markdown', 'gfm', 'commonmark', 'commonmark_x',
            'pdf', 'latex', 'docbook', 'docbook4', 'docbook5', 'jats', 'tei',
            'epub', 'epub2', 'epub3', 'fb2', 'rtf', 'txt', 'plain', 'json',
            'native', 'icml', 'opml', 'org', 'textile', 'mediawiki', 'dokuwiki',
            'haddock', 'man', 'ms', 'asciidoc', 'rst', 'docx', 'odt'
        ],
        'fb2': [
            'html', 'html5', 'xhtml', 'markdown', 'gfm', 'commonmark', 'commonmark_x',
            'pdf', 'latex', 'docbook', 'docbook4', 'docbook5', 'jats', 'tei',
            'epub', 'epub2', 'epub3', 'mobi', 'rtf', 'txt', 'plain', 'json',
            'native', 'icml', 'opml', 'org', 'textile', 'mediawiki', 'dokuwiki',
            'haddock', 'man', 'ms', 'asciidoc', 'rst', 'docx', 'odt'
        ],
        
        # Technical documentation formats
        'asciidoc': [
            'html', 'html5', 'xhtml', 'markdown', 'gfm', 'commonmark', 'commonmark_x',
            'pdf', 'latex', 'docbook', 'docbook4', 'docbook5', 'jats', 'tei',
            'epub', 'epub2', 'epub3', 'mobi', 'fb2', 'rtf', 'txt', 'plain',
            'json', 'native', 'icml', 'opml', 'org', 'textile', 'mediawiki',
            'dokuwiki', 'haddock', 'man', 'ms', 'rst', 'docx', 'odt'
        ],
        'rst': [
            'html', 'html5', 'xhtml', 'markdown', 'gfm', 'commonmark', 'commonmark_x',
            'pdf', 'latex', 'docbook', 'docbook4', 'docbook5', 'jats', 'tei',
            'epub', 'epub2', 'epub3', 'mobi', 'fb2', 'rtf', 'txt', 'plain',
            'json', 'native', 'icml', 'opml', 'org', 'textile', 'mediawiki',
            'dokuwiki', 'haddock', 'man', 'ms', 'asciidoc', 'docx', 'odt'
        ],
        'org': [
            'html', 'html5', 'xhtml', 'markdown', 'gfm', 'commonmark', 'commonmark_x',
            'pdf', 'latex', 'docbook', 'docbook4', 'docbook5', 'jats', 'tei',
            'epub', 'epub2', 'epub3', 'mobi', 'fb2', 'rtf', 'txt', 'plain',
            'json', 'native', 'icml', 'opml', 'textile', 'mediawiki', 'dokuwiki',
            'haddock', 'man', 'ms', 'asciidoc', 'rst', 'docx', 'odt'
        ],
        
        # Wiki formats
        'mediawiki': [
            'html', 'html5', 'xhtml', 'markdown', 'gfm', 'commonmark', 'commonmark_x',
            'pdf', 'latex', 'docbook', 'docbook4', 'docbook5', 'jats', 'tei',
            'epub', 'epub2', 'epub3', 'mobi', 'fb2', 'rtf', 'txt', 'plain',
            'json', 'native', 'icml', 'opml', 'org', 'textile', 'dokuwiki',
            'haddock', 'man', 'ms', 'asciidoc', 'rst', 'docx', 'odt'
        ],
        'dokuwiki': [
            'html', 'html5', 'xhtml', 'markdown', 'gfm', 'commonmark', 'commonmark_x',
            'pdf', 'latex', 'docbook', 'docbook4', 'docbook5', 'jats', 'tei',
            'epub', 'epub2', 'epub3', 'mobi', 'fb2', 'rtf', 'txt', 'plain',
            'json', 'native', 'icml', 'opml', 'org', 'textile', 'mediawiki',
            'haddock', 'man', 'ms', 'asciidoc', 'rst', 'docx', 'odt'
        ],
        
        # Plain text
        'txt': [
            'html', 'html5', 'xhtml', 'markdown', 'gfm', 'commonmark', 'commonmark_x',
            'pdf', 'latex', 'docbook', 'docbook4', 'docbook5', 'jats', 'tei',
            'epub', 'epub2', 'epub3', 'mobi', 'fb2', 'rtf', 'plain', 'json',
            'native', 'icml', 'opml', 'org', 'textile', 'mediawiki', 'dokuwiki',
            'haddock', 'man', 'ms', 'asciidoc', 'rst', 'docx', 'odt'
        ],
        'plain': [
            'html', 'html5', 'xhtml', 'markdown', 'gfm', 'commonmark', 'commonmark_x',
            'pdf', 'latex', 'docbook', 'docbook4', 'docbook5', 'jats', 'tei',
            'epub', 'epub2', 'epub3', 'mobi', 'fb2', 'rtf', 'txt', 'json',
            'native', 'icml', 'opml', 'org', 'textile', 'mediawiki', 'dokuwiki',
            'haddock', 'man', 'ms', 'asciidoc', 'rst', 'docx', 'odt'
        ],
        
        # LaTeX
        'latex': [
            'html', 'html5', 'xhtml', 'markdown', 'gfm', 'commonmark', 'commonmark_x',
            'pdf', 'docbook', 'docbook4', 'docbook5', 'jats', 'tei', 'epub',
            'epub2', 'epub3', 'mobi', 'fb2', 'rtf', 'txt', 'plain', 'json',
            'native', 'icml', 'opml', 'org', 'textile', 'mediawiki', 'dokuwiki',
            'haddock', 'man', 'ms', 'asciidoc', 'rst', 'docx', 'odt'
        ],
        'tex': [
            'html', 'html5', 'xhtml', 'markdown', 'gfm', 'commonmark', 'commonmark_x',
            'pdf', 'docbook', 'docbook4', 'docbook5', 'jats', 'tei', 'epub',
            'epub2', 'epub3', 'mobi', 'fb2', 'rtf', 'txt', 'plain', 'json',
            'native', 'icml', 'opml', 'org', 'textile', 'mediawiki', 'dokuwiki',
            'haddock', 'man', 'ms', 'asciidoc', 'rst', 'docx', 'odt'
        ]
    }
    
    # Get all supported output formats for the input formats
    all_supported_formats = set()
    
    for input_format in input_formats:
        if input_format in format_compatibility:
            all_supported_formats.update(format_compatibility[input_format])
    
    # If no specific compatibility found, return common formats
    if not all_supported_formats:
        all_supported_formats = {
            'html', 'html5', 'xhtml', 'markdown', 'gfm', 'commonmark', 'commonmark_x',
            'pdf', 'latex', 'docbook', 'docbook4', 'docbook5', 'jats', 'tei',
            'epub', 'epub2', 'epub3', 'mobi', 'fb2', 'rtf', 'txt', 'plain',
            'json', 'native', 'icml', 'opml', 'org', 'textile', 'mediawiki',
            'dokuwiki', 'haddock', 'man', 'ms', 'asciidoc', 'rst', 'docx', 'odt',
            'revealjs', 'beamer', 's5', 'slideous', 'dzslides', 'slidy'
        }
    
    # Sort formats for better presentation
    return sorted(list(all_supported_formats))

@app.route('/supported-formats', methods=['POST'])
def get_supported_formats():
    """Get supported output formats for the given input files"""
    try:
        if 'files' not in request.files:
            return jsonify({'error': 'No files provided'}), 400
        
        files = request.files.getlist('files')
        
        if not files or all(file.filename == '' for file in files):
            return jsonify({'error': 'No files selected'}), 400
        
        # Get input formats from uploaded files
        input_formats = set()
        for file in files:
            if file and file.filename:
                input_format = get_input_format(file.filename)
                input_formats.add(input_format)
        
        # Get supported output formats for these input formats
        supported_formats = get_supported_output_formats(input_formats)
        
        return jsonify({
            'input_formats': list(input_formats),
            'supported_output_formats': supported_formats,
            'message': f'Found {len(supported_formats)} supported output formats for {len(input_formats)} input format(s)'
        }), 200
        
    except Exception as e:
        logger.error(f"Error getting supported formats: {str(e)}")
        return jsonify({'error': f'Server error: {str(e)}'}), 500

@app.route('/all-formats', methods=['GET'])
def get_all_formats():
    """Get all supported input and output formats"""
    try:
        # Get all supported input formats
        input_formats = list(ALLOWED_EXTENSIONS)
        
        # Get all possible output formats
        all_output_formats = set()
        for input_format in ALLOWED_EXTENSIONS:
            input_format_name = get_input_format(f"test.{input_format}")
            supported = get_supported_output_formats([input_format_name])
            all_output_formats.update(supported)
        
        # Organize output formats by category
        format_categories = {
            'Web Formats': ['html', 'html5', 'xhtml'],
            'Document Formats': ['pdf', 'docx', 'odt', 'rtf'],
            'Markup Formats': ['markdown', 'gfm', 'commonmark', 'commonmark_x', 'markua'],
            'E-book Formats': ['epub', 'epub2', 'epub3', 'mobi', 'fb2'],
            'Technical Documentation': ['asciidoc', 'rst', 'org', 'textile', 'mediawiki', 'dokuwiki'],
            'Presentation Formats': ['revealjs', 'beamer', 's5', 'slideous', 'dzslides', 'slidy'],
            'Structured Formats': ['docbook', 'docbook4', 'docbook5', 'jats', 'tei', 'icml'],
            'Plain Text': ['txt', 'plain'],
            'Other Formats': ['json', 'native', 'opml', 'haddock', 'man', 'ms', 'latex']
        }
        
        organized_formats = {}
        for category, formats in format_categories.items():
            organized_formats[category] = [f for f in formats if f in all_output_formats]
        
        return jsonify({
            'input_formats': sorted(input_formats),
            'output_formats': sorted(list(all_output_formats)),
            'organized_output_formats': organized_formats,
            'message': f'Found {len(input_formats)} input formats and {len(all_output_formats)} output formats'
        }), 200
        
    except Exception as e:
        logger.error(f"Error getting all formats: {str(e)}")
        return jsonify({'error': f'Server error: {str(e)}'}), 500

if __name__ == '__main__':
    port = int(os.environ.get('PORT', 5000))
    app.run(use_reloader=False, debug=False, host='0.0.0.0', port=port)